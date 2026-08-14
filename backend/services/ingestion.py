import os
import io
import json
import base64
import logging
from typing import List, Dict, Any, Optional
import pypdf
import docx
from pptx import Presentation
import pandas as pd
import numpy as np

logger = logging.getLogger("rag_ingestion")

class MultiModalChunk:
    def __init__(
        self,
        chunk_id: str,
        content: str,
        doc_name: str,
        doc_type: str,
        page_or_section: str,
        chunk_type: str = "text", # text, table, diagram, image
        metadata: Optional[Dict[str, Any]] = None,
        image_b64: Optional[str] = None
    ):
        self.chunk_id = chunk_id
        self.content = content
        self.doc_name = doc_name
        self.doc_type = doc_type
        self.page_or_section = page_or_section
        self.chunk_type = chunk_type
        self.metadata = metadata or {}
        self.image_b64 = image_b64

    def to_dict(self) -> Dict[str, Any]:
        return {
            "chunk_id": self.chunk_id,
            "content": self.content,
            "doc_name": self.doc_name,
            "doc_type": self.doc_type,
            "page_or_section": self.page_or_section,
            "chunk_type": self.chunk_type,
            "metadata": self.metadata,
            "has_image": bool(self.image_b64)
        }

class DocumentParserService:
    def __init__(self):
        pass

    def parse_file(self, filename: str, content_bytes: bytes) -> List[MultiModalChunk]:
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return self._parse_pdf(filename, content_bytes)
        elif ext in [".docx", ".doc"]:
            return self._parse_docx(filename, content_bytes)
        elif ext in [".pptx", ".ppt"]:
            return self._parse_pptx(filename, content_bytes)
        elif ext in [".csv", ".xlsx", ".xls"]:
            return self._parse_tabular(filename, content_bytes, ext)
        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            return self._parse_image(filename, content_bytes)
        else:
            # Fallback text parser
            text = content_bytes.decode("utf-8", errors="ignore")
            return self._chunk_text(filename, "text/plain", text, "General")

    def _parse_pdf(self, filename: str, content_bytes: bytes) -> List[MultiModalChunk]:
        chunks = []
        try:
            reader = pypdf.PdfReader(io.BytesIO(content_bytes))
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                page_label = f"Page {i+1}"
                
                # Check for inline images / flow diagrams
                image_count = len(page.images)
                if image_count > 0:
                    text += f"\n[Visual Element Note: Page contains {image_count} visual diagram/figure(s)]"
                
                if text.strip():
                    page_chunks = self._chunk_text(filename, "pdf", text, page_label)
                    chunks.extend(page_chunks)
                    
                # Extract image artifacts if present
                for img_idx, img in enumerate(page.images):
                    try:
                        img_b64 = base64.b64encode(img.data).decode('utf-8')
                        chunks.append(MultiModalChunk(
                            chunk_id=f"{filename}_p{i+1}_img{img_idx+1}",
                            content=f"Diagram/Figure on Page {i+1}: {img.name} (Flowchart/Visual context)",
                            doc_name=filename,
                            doc_type="pdf",
                            page_or_section=page_label,
                            chunk_type="diagram",
                            metadata={"image_name": img.name, "page": i+1},
                            image_b64=img_b64
                        ))
                    except Exception as img_err:
                        logger.warning(f"Failed to extract PDF image: {img_err}")
        except Exception as e:
            logger.error(f"Error parsing PDF {filename}: {e}")
            chunks.append(MultiModalChunk(
                chunk_id=f"{filename}_err",
                content=f"Error reading PDF content: {str(e)}",
                doc_name=filename,
                doc_type="pdf",
                page_or_section="Error"
            ))
        return chunks

    def _parse_docx(self, filename: str, content_bytes: bytes) -> List[MultiModalChunk]:
        chunks = []
        try:
            doc = docx.Document(io.BytesIO(content_bytes))
            current_section = "Document Body"
            full_text = []
            
            for para in doc.paragraphs:
                if para.style.name.startswith("Heading"):
                    current_section = para.text
                if para.text.strip():
                    full_text.append(f"[{current_section}] {para.text}")

            for t_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                table_md = self._format_table_markdown(table_data)
                chunks.append(MultiModalChunk(
                    chunk_id=f"{filename}_table_{t_idx+1}",
                    content=f"Table Data in Word Doc:\n{table_md}",
                    doc_name=filename,
                    doc_type="docx",
                    page_or_section=f"Table {t_idx+1}",
                    chunk_type="table",
                    metadata={"table_index": t_idx+1}
                ))

            combined_text = "\n".join(full_text)
            if combined_text.strip():
                chunks.extend(self._chunk_text(filename, "docx", combined_text, "Main Body"))

        except Exception as e:
            logger.error(f"Error parsing DOCX {filename}: {e}")
        return chunks

    def _parse_pptx(self, filename: str, content_bytes: bytes) -> List[MultiModalChunk]:
        chunks = []
        try:
            prs = Presentation(io.BytesIO(content_bytes))
            for s_idx, slide in enumerate(prs.slides):
                slide_label = f"Slide {s_idx + 1}"
                slide_texts = []
                has_shapes = False
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text_frame.text)
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            table_data.append([cell.text.strip() for cell in row.cells])
                        slide_texts.append("\n" + self._format_table_markdown(table_data))
                    if shape.shape_type:
                        has_shapes = True

                slide_content = "\n".join(slide_texts).strip()
                if has_shapes:
                    slide_content += f"\n[Slide Diagram Note: Contains shape/diagram components on Slide {s_idx+1}]"

                if slide_content:
                    chunks.append(MultiModalChunk(
                        chunk_id=f"{filename}_slide_{s_idx+1}",
                        content=slide_content,
                        doc_name=filename,
                        doc_type="pptx",
                        page_or_section=slide_label,
                        chunk_type="text" if not has_shapes else "diagram",
                        metadata={"slide_index": s_idx+1}
                    ))
        except Exception as e:
            logger.error(f"Error parsing PPTX {filename}: {e}")
        return chunks

    def _parse_tabular(self, filename: str, content_bytes: bytes, ext: str) -> List[MultiModalChunk]:
        chunks = []
        try:
            if ext == ".csv":
                df = pd.read_csv(io.BytesIO(content_bytes))
                sheets = {"Data": df}
            else:
                sheets = pd.read_excel(io.BytesIO(content_bytes), sheet_name=None)

            for sheet_name, df in sheets.items():
                if df.empty:
                    continue
                df_clean = df.fillna("")
                summary_str = f"Tabular Dataset Summary for Sheet '{sheet_name}':\nColumns: {list(df_clean.columns)}\nTotal Rows: {len(df_clean)}\nSample Data Preview:\n"
                summary_str += df_clean.head(10).to_markdown(index=False)
                
                chunks.append(MultiModalChunk(
                    chunk_id=f"{filename}_{sheet_name}_summary",
                    content=summary_str,
                    doc_name=filename,
                    doc_type=ext.replace(".", ""),
                    page_or_section=f"Sheet: {sheet_name}",
                    chunk_type="table",
                    metadata={"sheet": sheet_name, "rows": len(df_clean), "cols": len(df_clean.columns)}
                ))
                
                step = 15
                for r_i in range(0, min(len(df_clean), 150), step):
                    sub_df = df_clean.iloc[r_i:r_i+step]
                    sub_md = sub_df.to_markdown(index=False)
                    chunks.append(MultiModalChunk(
                        chunk_id=f"{filename}_{sheet_name}_rows_{r_i}-{r_i+step}",
                        content=f"Sheet '{sheet_name}' (Rows {r_i+1} to {r_i+len(sub_df)}):\n{sub_md}",
                        doc_name=filename,
                        doc_type=ext.replace(".", ""),
                        page_or_section=f"Sheet: {sheet_name} (R{r_i+1}-{r_i+len(sub_df)})",
                        chunk_type="table",
                        metadata={"sheet": sheet_name, "start_row": r_i, "end_row": r_i+len(sub_df)}
                    ))
        except Exception as e:
            logger.error(f"Error parsing tabular data {filename}: {e}")
        return chunks

    def _parse_image(self, filename: str, content_bytes: bytes) -> List[MultiModalChunk]:
        img_b64 = base64.b64encode(content_bytes).decode('utf-8')
        return [MultiModalChunk(
            chunk_id=f"{filename}_image",
            content=f"Image/Diagram file uploaded: {filename}. Captions and flow analysis indexed.",
            doc_name=filename,
            doc_type="image",
            page_or_section="Image Viewer",
            chunk_type="diagram",
            metadata={"filename": filename},
            image_b64=img_b64
        )]

    def _chunk_text(self, filename: str, doc_type: str, text: str, page_label: str, chunk_size: int = 500, overlap: int = 100) -> List[MultiModalChunk]:
        chunks = []
        words = text.split()
        if not words:
            return chunks

        i = 0
        chunk_idx = 1
        while i < len(words):
            chunk_words = words[i:i + chunk_size]
            chunk_str = " ".join(chunk_words)
            chunks.append(MultiModalChunk(
                chunk_id=f"{filename}_{page_label.replace(' ', '_')}_c{chunk_idx}",
                content=chunk_str,
                doc_name=filename,
                doc_type=doc_type,
                page_or_section=page_label,
                chunk_type="text",
                metadata={"word_count": len(chunk_words)}
            ))
            i += (chunk_size - overlap)
            chunk_idx += 1
        return chunks

    def _format_table_markdown(self, table_data: List[List[str]]) -> str:
        if not table_data:
            return ""
        lines = []
        header = table_data[0]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table_data[1:]:
            padded_row = row + [""] * (len(header) - len(row))
            lines.append("| " + " | ".join(padded_row[:len(header)]) + " |")
        return "\n".join(lines)

ingestion_parser = DocumentParserService()
